"""Pipeline de ingestao: loaders por formato -> limpeza -> chunking -> metadados.

Um processador por formato dos 8 documentos do EduNova (docs/raw/), metadados
(arquivo, tema, data_atualizacao, status) herdados de docs/catalogo.csv.
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path

from langchain_community.document_loaders import (
    BSHTMLLoader,
    Docx2txtLoader,
    PyPDFLoader,
    TextLoader,
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parent.parent
RAW_DIR = BASE_DIR / "docs" / "raw"
CATALOGO_PATH = BASE_DIR / "docs" / "catalogo.csv"

CHUNK_SIZE = 800
CHUNK_OVERLAP = 120

# Remove apenas caracteres de controle; preserva \t (0x09) e \n (0x0a).
_CONTROL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
_ESPACOS_RE = re.compile(r"[ \t]+")
_LINHAS_EM_BRANCO_RE = re.compile(r"\n{3,}")

_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP
)


def carregar_catalogo() -> dict[str, dict[str, str]]:
    with CATALOGO_PATH.open(encoding="utf-8") as f:
        return {row["arquivo"]: row for row in csv.DictReader(f)}


def limpar_texto(texto: str) -> str:
    texto = _CONTROL_CHARS_RE.sub("", texto)
    texto = _ESPACOS_RE.sub(" ", texto)
    texto = _LINHAS_EM_BRANCO_RE.sub("\n\n", texto)
    return texto.strip()


def _metadados_base(arquivo: str, catalogo: dict[str, dict[str, str]]) -> dict[str, str]:
    info = catalogo[arquivo]
    return {
        "arquivo": arquivo,
        "tema": info["tema"],
        "data_atualizacao": info["data_atualizacao"],
        # `obsoleto` marca a versao antiga de um documento ja revisado: continua
        # indexada (auditoria/curadoria) mas o retriever nao a entrega ao LLM.
        "status": info["status"],
    }


def _chunk_texto_corrido(texto: str, metadados: dict[str, str]) -> list[Document]:
    texto = limpar_texto(texto)
    if not texto:
        return []
    return [
        Document(page_content=parte, metadata={**metadados, "chunk_index": i})
        for i, parte in enumerate(_splitter.split_text(texto))
    ]


def _documento_por_registro(
    registros: list[dict],
    metadados: dict[str, str],
    campo_dono: str | None = None,
    campos_extra: list[str] | None = None,
) -> list[Document]:
    documentos = []
    for i, registro in enumerate(registros):
        conteudo = "\n".join(f"{campo}: {valor}" for campo, valor in registro.items())
        meta = {**metadados, "chunk_index": i}
        if campo_dono and registro.get(campo_dono):
            meta["dono"] = registro[campo_dono]
        for campo in campos_extra or []:
            if registro.get(campo) is not None:
                meta[campo] = registro[campo]
        documentos.append(Document(page_content=limpar_texto(conteudo), metadata=meta))
    return documentos


def processar_pdf(caminho: Path, metadados: dict[str, str]) -> list[Document]:
    paginas = PyPDFLoader(str(caminho)).load()
    texto = "\n\n".join(p.page_content for p in paginas)
    return _chunk_texto_corrido(texto, metadados)


def processar_docx(caminho: Path, metadados: dict[str, str]) -> list[Document]:
    texto = "\n\n".join(d.page_content for d in Docx2txtLoader(str(caminho)).load())
    return _chunk_texto_corrido(texto, metadados)


def processar_markdown(caminho: Path, metadados: dict[str, str]) -> list[Document]:
    texto = TextLoader(str(caminho), encoding="utf-8").load()[0].page_content
    return _chunk_texto_corrido(texto, metadados)


def processar_html(caminho: Path, metadados: dict[str, str]) -> list[Document]:
    texto = BSHTMLLoader(str(caminho), open_encoding="utf-8").load()[0].page_content
    return _chunk_texto_corrido(texto, metadados)


def processar_pptx(caminho: Path, metadados: dict[str, str]) -> list[Document]:
    apresentacao = Presentation(str(caminho))
    slides_texto = []
    for slide in apresentacao.slides:
        linhas = [
            paragrafo.text
            for shape in slide.shapes
            if shape.has_text_frame
            for paragrafo in shape.text_frame.paragraphs
            if paragrafo.text.strip()
        ]
        if linhas:
            slides_texto.append("\n".join(linhas))
    texto = "\n\n".join(slides_texto)
    return _chunk_texto_corrido(texto, metadados)


def processar_xlsx(caminho: Path, metadados: dict[str, str]) -> list[Document]:
    planilha = load_workbook(caminho, read_only=True, data_only=True).active
    linhas = list(planilha.iter_rows(values_only=True))
    if not linhas:
        return []
    cabecalho, *dados = linhas
    registros = [
        {coluna: valor for coluna, valor in zip(cabecalho, linha) if valor is not None}
        for linha in dados
    ]
    return _documento_por_registro(registros, metadados)


def processar_csv(
    caminho: Path,
    metadados: dict[str, str],
    campo_dono: str | None = None,
    campos_extra: list[str] | None = None,
) -> list[Document]:
    with caminho.open(encoding="utf-8") as f:
        registros = list(csv.DictReader(f))
    return _documento_por_registro(
        registros, metadados, campo_dono=campo_dono, campos_extra=campos_extra
    )


def processar_json(caminho: Path, metadados: dict[str, str]) -> list[Document]:
    with caminho.open(encoding="utf-8") as f:
        dados = json.load(f)
    # Assume um unico array de nivel superior (ex.: {"cursos": [...]})
    (registros,) = dados.values()
    return _documento_por_registro(registros, metadados)


_PROCESSADORES = {
    ".pdf": processar_pdf,
    ".docx": processar_docx,
    ".md": processar_markdown,
    ".html": processar_html,
    ".pptx": processar_pptx,
    ".xlsx": processar_xlsx,
    ".csv": processar_csv,
    ".json": processar_json,
}


# Documentos do corpus com dado pessoal identificavel. Marcar o "dono" de cada
# chunk permite ao no `autorizar` do grafo (Passo 5/7) filtrar por estudante
# identificado antes do `generate`. `matriculas_alunos.csv` tambem carrega
# `data_matricula`/`percentual_concluido` como metadado (nao so no texto do
# chunk) para o calculo determinístico de reembolso personalizado (Passo 8).
_ARQUIVOS_COM_DONO = {
    "certificados_emitidos.csv": {"campo_dono": "aluno"},
    "matriculas_alunos.csv": {
        "campo_dono": "aluno",
        "campos_extra": ["data_matricula", "percentual_concluido"],
    },
}


def processar_documentos() -> list[Document]:
    catalogo = carregar_catalogo()
    documentos: list[Document] = []
    for arquivo in catalogo:
        caminho = RAW_DIR / arquivo
        processador = _PROCESSADORES[caminho.suffix.lower()]
        metadados = _metadados_base(arquivo, catalogo)
        config_dono = _ARQUIVOS_COM_DONO.get(arquivo)
        if config_dono:
            documentos.extend(processador(caminho, metadados, **config_dono))
        else:
            documentos.extend(processador(caminho, metadados))
    return documentos


if __name__ == "__main__":
    catalogo = carregar_catalogo()
    docs = processar_documentos()
    print(f"{len(docs)} chunks gerados a partir de {len(catalogo)} documentos.\n")
    por_arquivo: dict[str, int] = {}
    for d in docs:
        por_arquivo[d.metadata["arquivo"]] = por_arquivo.get(d.metadata["arquivo"], 0) + 1
    for arquivo in catalogo:
        print(f"  {arquivo}: {por_arquivo.get(arquivo, 0)} chunk(s)")
